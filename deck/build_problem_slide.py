#!/usr/bin/env python3
"""
Builds deck/01-problem-framing.pptx — the one-slide problem framing for the
Aberdeen Advisors hackathon (Team 1, Transformation Roadmap Generator).

Every figure on the slide is derived from synthetic-data/. See FIGURES below
for the provenance of each number.

Run:  pip install python-pptx --break-system-packages
      python3 deck/build_problem_slide.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# FIGURES — provenance (all from synthetic-data/, currency is USD)
# ---------------------------------------------------------------------------
# 60 initiatives / 8 functions ....... rowcount + distinct `function` in data/initiatives.csv
# $95.3m approved .................... sum(total_budget) in data/initiatives.csv = 95,350,000
# $101.7m forecast ................... sum(forecast_at_completion)              = 101,743,000
# $84.3m claimed annual benefit ...... sum(annual_benefit_target)               = 84,270,000
# 77% Med/Low confidence ............. sum(annual_benefit_target) where value_confidence != High
#                                      = 65,220,000 / 84,270,000
# $18.8m spent / ~$0.5k banked ....... sum(actual_spend) and sum(benefit_actual) in
#                                      data/burn.csv and data/benefits.csv (blank from 2026-08)
# 95 dependencies / 70 cross-function  rowcount of data/dependencies.csv; cross = endpoints
#                                      whose `function` in initiatives.csv differs (70/95 = 74%)
# 3 violations, worst 716 days ....... data/dependency_conflicts.csv (INIT-039 -> INIT-031,
#                                      HR / People gating Finance, overlap_days = 716)
# 3 roles, peak 155% ................. data/resources.csv, over_allocated = TRUE:
#                                      Data Engineer, Supply Chain Analyst, Change Manager;
#                                      max utilisation_pct = 154.8 (Supply Chain Analyst, 2026-10)
# 164 of 259 milestones slipping ..... data/milestones.csv, slip_days > 0
# $895k enabler gating 28 / $38.6m ... INIT-001 total_budget 895,000, annual_benefit_target 70,000;
#                                      transitive successors in dependencies.csv = 28 initiatives,
#                                      sum(annual_benefit_target) of those = 38,580,000
# 3 scenarios / -25% capex / $8.3m ... data/scenarios.json: SC-01/02/03; capex_cap_usd falls
#                                      33,352,000 -> 25,014,000 (-8,338,000); SC-02 & SC-03
#                                      defer 6 initiatives worth $10.36m budget / $5.75m benefit
# 5 report-vs-data contradictions .... e.g. docs/status-reports/2026-07-pmo-monthly-report.md
#                                      RAG table totals 32/18/10 vs 31/17/12 in initiatives.csv
# 3 date formats / phantom row ....... docs/pmo-tracker-messy.csv (DD/MM/YY, ISO, MMM-YY;
#                                      INIT-002a exists nowhere else; INIT-011 & INIT-020 dup)
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x16, 0x28, 0x3C)
BODY = RGBColor(0x33, 0x3F, 0x4F)
ACCENT = RGBColor(0x1B, 0x6E, 0x8C)
MUTED = RGBColor(0x5A, 0x66, 0x74)
PANEL = RGBColor(0xF4, 0xF6, 0xF8)
RULE = RGBColor(0xD5, 0xDC, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE = "Transformation Roadmap Generator"
SUBTITLE = ("Turning 60 inherited initiatives, $95.3m of spend and 95 tangled dependencies "
            "into a sequenced, board-ready roadmap")

SECTIONS = [
    ("THE CHALLENGE", [
        "A new leader inherits **60 initiatives** across **8 functions** — no integrated roadmap, no agreed priorities",
        "**$95.3m** approved, forecasting **$101.7m**; **$84.3m** claimed annual benefit, **~$0.5k** banked",
        "**95** dependencies, **70** cross-function, in nobody's plan — **3** violated by up to **716 days**",
        "\"What happens first, and when does value land?\" takes **weeks** of PMO work — and arrives stale",
    ]),
    ("TARGET USER", [
        "**Primary: the Chief Transformation Officer**, first 90 days, owing the board a credible plan",
        "**Secondary: the PMO lead** who maintains it across **8** business cases and **3** monthly reports",
        "**Secondary: the CFO**, who consumes the value view and defends the **$84.3m** benefit baseline",
        "Not \"executives\": one named owner, a monthly board cycle, accountable for the number",
    ]),
    ("THE WORKFLOW", [
        "**In:** the plans that exist — business cases, PMO reports, steering minutes, the master tracker",
        "**Normalise:** initiatives, cost, dates, dependencies — **3** date formats, **6** budget formats, duplicates",
        "**Challenge:** the **716-day** violation, **5** report-vs-data contradictions, **1** phantom initiative",
        "**Out:** sequenced roadmap + dependency view → **3** scenarios compared → value-realization plan",
    ]),
    ("BUSINESS VALUE", [
        "**Weeks of PMO analysis → minutes**, rebuilt on every status report, not once a quarter",
        "Finds the **70** cross-function edges manual planning misses — an **$895k** enabler gates **$38.6m/yr**",
        "Flags **3** violations, **164** slipping milestones, **3** roles at **155%** — before the board commits",
        "Prices the real choice: **−25% capex (−$8.3m)** vs speed-to-value, with **$10.4m** of deferrals costed",
    ]),
]


def add_runs(p, text, size, color):
    """Mini-markup: **bold** segments render bold + accent-dark."""
    for i, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        r.font.size = Pt(size)
        r.font.name = "Calibri"
        bold = (i % 2 == 1)
        r.font.bold = bold
        r.font.color.rgb = NAVY if bold else color


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ---- title -------------------------------------------------------
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.23), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = TITLE
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = NAVY

    sb = slide.shapes.add_textbox(Inches(0.55), Inches(0.86), Inches(12.23), Inches(0.40))
    stf = sb.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_right = stf.margin_top = stf.margin_bottom = 0
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = SUBTITLE
    sr.font.size = Pt(14)
    sr.font.name = "Calibri"
    sr.font.color.rgb = MUTED

    # accent rule under the title block
    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(1.32), Inches(12.23), Pt(2.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False

    # ---- 2x2 quadrants -----------------------------------------------
    left_x, right_x = 0.55, 6.92
    top_y, bot_y = 1.62, 4.58
    w, h = 5.86, 2.78

    for idx, (heading, bullets) in enumerate(SECTIONS):
        x = left_x if idx % 2 == 0 else right_x
        y = top_y if idx < 2 else bot_y

        panel = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = RULE
        panel.line.width = Pt(0.75)
        panel.shadow.inherit = False

        bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Pt(4.5), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        bar.shadow.inherit = False

        box = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 0.14),
                                       Inches(w - 0.38), Inches(h - 0.24))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP

        hp = tf.paragraphs[0]
        hr = hp.add_run()
        hr.text = f"{idx + 1}.  {heading}"
        hr.font.size = Pt(14)
        hr.font.bold = True
        hr.font.name = "Calibri"
        hr.font.color.rgb = ACCENT
        hp.space_after = Pt(7)

        for b in bullets:
            bp = tf.add_paragraph()
            add_runs(bp, "•  " + b, 12.5, BODY)
            bp.space_after = Pt(6)
            bp.line_spacing = 0.96
            # hanging indent so wrapped lines align under the text, not the bullet
            bp._pPr.set("marL", str(Inches(0.20)))
            bp._pPr.set("indent", str(-Inches(0.20)))

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print("wrote", path)


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    build(os.path.join(here, "01-problem-framing.pptx"))
