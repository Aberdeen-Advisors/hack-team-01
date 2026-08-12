#!/usr/bin/env python3
"""
Builds deck/transformation-roadmap-generator-4slide.pptx, the SHORT version of
the Team 1 hackathon deck: four slides, plain English, heavily visual.

The long 8-slide deck (deck/build_deck.py, transformation-roadmap-generator.pptx)
stays exactly as it is. This file only adds a shorter companion. Palette, fonts
and the header helper are imported from build_deck.py so both decks look like the
same family; type here is bigger and word counts are much lower.

Style rules for this deck:
  * no em dashes or en dashes anywhere in the generated text (see check_dashes)
  * short words, low word count, graphics carry each slide

All figures are USD and trace to synthetic-data/. NOTE: the "$84.3m claimed vs
$507 banked" comparison was retired as wrong, never reintroduce it. The correct
like-for-like line is $84.3m promised vs a $39.6m annual run-rate scheduled by
end-2027 (47%).

Run:  pip install python-pptx --break-system-packages
      python3 deck/build_short_deck.py
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from build_deck import (  # noqa: E402
    NAVY, BODY, ACCENT, MUTED, PANEL, RULE, WHITE,
    new_slide, header, panel, RECT,
)

# Two extra tints mixed from the existing palette.
TEAL_LT = RGBColor(0x8F, 0xC2, 0xD2)
AMBER = RGBColor(0xC8, 0x7A, 0x1E)


# ---------------------------------------------------------------------------
# small drawing helpers (shapes only, no images, no network)
# ---------------------------------------------------------------------------

def text(slide, x, y, w, h, runs, align_top=True):
    """runs = list of (string, size_pt, bold, color, space_after_pt)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP if align_top else MSO_ANCHOR.MIDDLE
    for i, (s, size, bold, color, after) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Calibri"
        r.font.color.rgb = color
        p.space_after = Pt(after)
        p.line_spacing = 0.98
    return tf


def block(slide, x, y, w, h, fill, line=None):
    b = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    b.fill.solid()
    b.fill.fore_color.rgb = fill
    if line is None:
        b.line.fill.background()
    else:
        b.line.color.rgb = line
        b.line.width = Pt(0.75)
    b.shadow.inherit = False
    return b


def big_arrow(slide, x, y, w, h, color=ACCENT):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def strip(slide, x, y, w, h, heading, body, fill=NAVY,
          head_color=TEAL_LT, body_color=WHITE, body_size=12.5):
    block(slide, x, y, w, h, fill)
    runs = []
    if heading:
        runs.append((heading, 13, True, head_color, 5))
    runs.append((body, body_size, False, body_color, 0))
    text(slide, x + 0.28, y + 0.18, w - 0.56, h - 0.32, runs,
         align_top=bool(heading))


TINTS = [NAVY, ACCENT, TEAL_LT, AMBER]


# ---------------------------------------------------------------------------
# SLIDE 1 = title
# ---------------------------------------------------------------------------

def slide_1(prs):
    slide = new_slide(prs)

    # full-bleed navy field with an accent band
    block(slide, 0, 0, 13.333, 7.5, NAVY)
    block(slide, 0, 0, 0.34, 7.5, ACCENT)

    # decorative sequenced blocks, echoing the roadmap idea
    for i in range(14):
        block(slide, 1.10 + i * 0.60, 1.42, 0.44, 0.26, TINTS[i % 4])
    block(slide, 1.10, 1.86, 8.24, 0.05, ACCENT)

    text(slide, 1.10, 2.50, 11.60, 1.30, [
        ("Transformation Roadmap Generator", 42, True, WHITE, 0)])

    text(slide, 1.10, 3.86, 10.4, 0.70, [
        ("Turning 60 competing initiatives into one sequenced roadmap "
         "the board can agree on.", 19, False, TEAL_LT, 0)])

    block(slide, 1.10, 5.06, 3.10, 0.05, ACCENT)

    text(slide, 1.10, 5.40, 11.0, 1.20, [
        ("Team 1: Suraj Sehgal, Katie Bui, Jim Towler", 15, True, WHITE, 6),
        ("Hackathon 2026", 14, False, TEAL_LT, 0),
    ])


# ---------------------------------------------------------------------------
# SLIDE 2 = executive summary (problem / user / solution)
# ---------------------------------------------------------------------------

def slide_2(prs):
    slide = new_slide(prs)
    header(slide,
           "60 initiatives, no agreed order",
           "The problem, the person with the problem, and what we built for them.",
           eyebrow="EXECUTIVE SUMMARY")

    w, gap, y, h = 3.95, 0.19, 1.82, 4.24
    xs = [0.55 + i * (w + gap) for i in range(3)]

    # ---- panel 1: the problem, with a mess-to-order mini graphic ----------
    panel(slide, xs[0], y, w, h, fill=PANEL, bar=False)
    block(slide, xs[0], y, w, 0.16, AMBER)
    text(slide, xs[0] + 0.26, y + 0.40, w - 0.52, 0.34,
         [("THE PROBLEM", 13, True, AMBER, 0)])

    scatter = [(0.18, 0.10), (0.86, 0.02), (1.52, 0.20), (0.34, 0.56),
               (1.02, 0.62), (1.66, 0.48), (0.10, 1.02), (0.78, 1.10),
               (1.46, 0.96), (0.50, 1.48), (1.20, 1.42), (1.72, 1.54)]
    gx, gy = xs[0] + 0.26, y + 0.92
    for i, (dx, dy) in enumerate(scatter):
        block(slide, gx + dx, gy + dy, 0.42, 0.24, TINTS[i % 4])
    big_arrow(slide, gx + 2.28, gy + 0.70, 0.42, 0.34)
    for i in range(12):
        block(slide, gx + 2.70 + (i % 2) * 0.46, gy + 0.06 + (i // 2) * 0.32,
              0.42, 0.24, TINTS[i % 4])

    text(slide, xs[0] + 0.26, y + 3.00, w - 0.52, 1.10, [
        ("60 projects across technology, operations, growth and cost. "
         "No single plan. No agreement on what comes first.", 12.5, False, BODY, 0)])

    # ---- panel 2: the target user ----------------------------------------
    panel(slide, xs[1], y, w, h, fill=PANEL, bar=False)
    block(slide, xs[1], y, w, 0.16, ACCENT)
    text(slide, xs[1] + 0.26, y + 0.40, w - 0.52, 0.34,
         [("THE USER", 13, True, ACCENT, 0)])

    text(slide, xs[1] + 0.26, y + 0.92, w - 0.52, 1.10, [
        ("Chief Transformation Officer", 20, True, NAVY, 4),
        ("in their first 90 days.", 15, False, MUTED, 0),
    ])
    # two supporting-cast chips
    for i, who in enumerate(["Programme office lead", "Finance chief (CFO)"]):
        cy = y + 2.14 + i * 0.52
        block(slide, xs[1] + 0.26, cy, w - 0.52, 0.42, WHITE, line=RULE)
        text(slide, xs[1] + 0.44, cy + 0.11, w - 0.88, 0.30,
             [(who, 12.5, False, BODY, 0)])

    text(slide, xs[1] + 0.26, y + 3.32, w - 0.52, 0.80, [
        ("The board wants three answers: what happens first, what depends on what, "
         "and when the money shows up in the accounts.", 12.5, False, BODY, 0)])

    # ---- panel 3: our solution -------------------------------------------
    panel(slide, xs[2], y, w, h, fill=NAVY, bar=False)
    block(slide, xs[2], y, w, 0.16, TEAL_LT)
    text(slide, xs[2] + 0.26, y + 0.40, w - 0.52, 0.34,
         [("OUR SOLUTION", 13, True, TEAL_LT, 0)])

    text(slide, xs[2] + 0.26, y + 0.92, w - 0.52, 2.60, [
        ("One web page that reads the plans a client already has, puts every project "
         "into one consistent list, works out the order that respects what depends on "
         "what and the limits on money and people, then shows the roadmap, the links "
         "between projects and the benefit curve.", 14, False, WHITE, 8),
        ("Weeks of analysis in minutes.", 14, True, TEAL_LT, 0),
    ])


# ---------------------------------------------------------------------------
# SLIDE 3 = the visual workflow
# ---------------------------------------------------------------------------

def slide_3(prs):
    slide = new_slide(prs)
    header(slide,
           "Plans in, sequenced roadmap out",
           "Four steps, start to finish, in minutes.",
           eyebrow="HOW IT WORKS")

    steps = [
        ("1", "Read the plans", "Spreadsheets, business cases, status reports"),
        ("2", "One list", "Every project described the same way"),
        ("3", "Work out the order", "Dependencies, budget and people limits"),
        ("4", "Show the picture", "Roadmap, project links, benefit curve"),
    ]
    w, gap, y, h = 2.86, 0.24, 1.90, 2.60
    for i, (n, title, label) in enumerate(steps):
        x = 0.55 + i * (w + gap)
        fill = NAVY if i == 3 else PANEL
        panel(slide, x, y, w, h, fill=fill, bar=False)
        block(slide, x, y, w, 0.16, TINTS[i % 4])
        text(slide, x + 0.24, y + 0.42, w - 0.48, 0.90,
             [(n, 54, True, TEAL_LT if i == 3 else TINTS[i % 4], 0)])
        text(slide, x + 0.24, y + 1.42, w - 0.48, 1.00, [
            (title, 17, True, WHITE if i == 3 else NAVY, 6),
            (label, 12, False, TEAL_LT if i == 3 else MUTED, 0),
        ])
        if i < 3:
            big_arrow(slide, x + w + 0.03, y + 1.06, 0.18, 0.48)

    # ---- demo moment -------------------------------------------------------
    block(slide, 0.55, 4.80, 12.23, 2.00, PANEL, line=RULE)
    block(slide, 0.55, 4.80, 12.23, 0.16, AMBER)
    text(slide, 0.85, 5.14, 5.60, 1.40, [
        ("Ranked on return alone", 15, True, NAVY, 6),
        ("An $895,000 platform project looks like the worst investment "
         "in the portfolio, so it goes last.", 13, False, BODY, 0),
    ])
    big_arrow(slide, 6.66, 5.56, 0.62, 0.46, AMBER)
    text(slide, 7.56, 5.14, 5.00, 1.40, [
        ("It unlocks 28 other projects", 15, True, AMBER, 6),
        ("worth $38.6m a year. The tool moves it to the front. Same inputs, "
         "same answer, every time.", 13, False, BODY, 0),
    ])


# ---------------------------------------------------------------------------
# SLIDE 4 = the value, and what happens next
# ---------------------------------------------------------------------------

def tile(slide, x, y, w, h, value, caption, accent=ACCENT):
    """Big-number tile: coloured cap, huge figure, plain-English caption."""
    panel(slide, x, y, w, h, fill=PANEL, bar=False)
    block(slide, x, y, w, 0.16, accent)
    text(slide, x + 0.24, y + 0.44, w - 0.48, 0.85, [(value, 44, True, NAVY, 0)])
    text(slide, x + 0.24, y + 1.34, w - 0.48, h - 1.50,
         [(caption, 12.5, False, BODY, 0)])


def slide_4(prs):
    slide = new_slide(prs)
    header(slide,
           "Where the value actually lands",
           "The promise, the plan, and the gap between them.",
           eyebrow="THE VALUE")

    tiles = [
        ("47%", "of the $84.3m promised is actually scheduled: a $39.6m a year "
                "run-rate by the end of 2027.", ACCENT),
        ("$18.0m", "of benefit arrives inside the two year window, against $89.8m "
                   "of planned spend in the same window.", NAVY),
        ("16 of 60", "projects deliver nothing at all inside the two year window.",
         AMBER),
    ]
    w, gap = 3.95, 0.19
    for i, (v, c, a) in enumerate(tiles):
        tile(slide, 0.55 + i * (w + gap), 1.82, w, 2.26, v, c, accent=a)

    # ---- what happens next -------------------------------------------------
    block(slide, 0.55, 4.32, 12.23, 1.86, PANEL, line=RULE)
    block(slide, 0.55, 4.32, 12.23, 0.16, ACCENT)
    nexts = [
        (0.85, 3.60, "What's next",
         "A 4 to 6 week pilot on one client's real portfolio."),
        (4.85, 4.00, "What it needs",
         "A read-only export of their project tracker, a part-time analyst, "
         "a named executive sponsor."),
        (9.20, 3.35, "How it's judged",
         "Time to a board-ready roadmap, and clashes caught before money "
         "is committed."),
    ]
    for x, tw, head, body in nexts:
        text(slide, x, 4.64, tw, 1.40, [
            (head, 14, True, NAVY, 6),
            (body, 12.5, False, BODY, 0),
        ])

    strip(slide, 0.55, 6.38, 12.23, 0.72, "",
          "Weeks of programme office analysis in minutes, and it all runs in the "
          "client's own browser, with no data leaving their environment.")


# ---------------------------------------------------------------------------

BUILDERS = [slide_1, slide_2, slide_3, slide_4]

BANNED = {"—": "em dash", "–": "en dash"}


def check_dashes(prs):
    """Hard style rule: no em dashes or en dashes anywhere in the deck."""
    bad = []
    for n, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for ch, name in BANNED.items():
                if ch in shape.text_frame.text:
                    bad.append(f"slide {n}: {name} in {shape.text_frame.text[:60]!r}")
    if bad:
        raise SystemExit("banned characters found:\n  " + "\n  ".join(bad))


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in BUILDERS:
        fn(prs)
    check_dashes(prs)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"wrote {path} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    build(os.path.join(here, "transformation-roadmap-generator-4slide.pptx"))
