#!/usr/bin/env python3
"""
Builds deck/transformation-roadmap-generator.pptx — the full Team 1 hackathon
deck (slides 1-8) for the Transformation Roadmap Generator.

Slide 1 is regenerated from the SECTIONS/TITLE/SUBTITLE data in
build_problem_slide.py so the two files can never drift; slides 2-8 are built
here in the same visual language (Calibri, navy/teal, panel cards).

All figures trace to synthetic-data/ (USD). See the FIGURES block in
build_problem_slide.py for provenance. NOTE: the "$84.3m claimed vs $507
banked" comparison was retired as wrong — never reintroduce it. The correct
like-for-like line is $84.3m promised vs a $39.6m annual run-rate scheduled by
end-2027 (47%).

Run:  pip install python-pptx --break-system-packages
      python3 deck/build_deck.py
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from build_problem_slide import (  # noqa: E402
    SECTIONS as PROBLEM_SECTIONS,
    TITLE as DECK_TITLE,
    SUBTITLE as PROBLEM_SUBTITLE,
    NAVY, BODY, ACCENT, MUTED, PANEL, RULE, WHITE,
    add_runs,
)

RECT = MSO_SHAPE.RECTANGLE


# ---------------------------------------------------------------------------
# shared layout helpers (same geometry + type ramp as slide 1)
# ---------------------------------------------------------------------------

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def header(slide, title, subtitle, eyebrow=None):
    """Title + subtitle + accent rule, matching slide 1's header block."""
    if eyebrow:
        eb = slide.shapes.add_textbox(Inches(0.55), Inches(0.20), Inches(12.23), Inches(0.26))
        etf = eb.text_frame
        etf.word_wrap = True
        etf.margin_left = etf.margin_right = etf.margin_top = etf.margin_bottom = 0
        er = etf.paragraphs[0].add_run()
        er.text = eyebrow
        er.font.size = Pt(11)
        er.font.bold = True
        er.font.name = "Calibri"
        er.font.color.rgb = ACCENT

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.44 if eyebrow else 0.30),
                                  Inches(12.23), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = NAVY

    sb = slide.shapes.add_textbox(Inches(0.55), Inches(0.98 if eyebrow else 0.86),
                                  Inches(12.23), Inches(0.40))
    stf = sb.text_frame
    stf.word_wrap = True
    stf.margin_left = stf.margin_right = stf.margin_top = stf.margin_bottom = 0
    sr = stf.paragraphs[0].add_run()
    sr.text = subtitle
    sr.font.size = Pt(14)
    sr.font.name = "Calibri"
    sr.font.color.rgb = MUTED

    rule = slide.shapes.add_shape(RECT, Inches(0.55), Inches(1.44 if eyebrow else 1.32),
                                  Inches(12.23), Pt(2.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    rule.shadow.inherit = False


def panel(slide, x, y, w, h, fill=PANEL, bar=True, bar_color=ACCENT):
    p = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.color.rgb = RULE
    p.line.width = Pt(0.75)
    p.shadow.inherit = False
    if bar:
        b = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Pt(4.5), Inches(h))
        b.fill.solid()
        b.fill.fore_color.rgb = bar_color
        b.line.fill.background()
        b.shadow.inherit = False
    return p


def card(slide, x, y, w, h, heading, bullets, number=None,
         head_size=14, body_size=12.5, fill=PANEL, bar_color=ACCENT,
         head_color=ACCENT, body_color=BODY):
    """Panel card with a heading and bulleted lines — slide 1's quadrant style."""
    panel(slide, x, y, w, h, fill=fill, bar_color=bar_color)
    box = slide.shapes.add_textbox(Inches(x + 0.20), Inches(y + 0.14),
                                   Inches(w - 0.38), Inches(h - 0.24))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    hp = tf.paragraphs[0]
    hr = hp.add_run()
    hr.text = f"{number}.  {heading}" if number else heading
    hr.font.size = Pt(head_size)
    hr.font.bold = True
    hr.font.name = "Calibri"
    hr.font.color.rgb = head_color
    hp.space_after = Pt(7)

    for b in bullets:
        bp = tf.add_paragraph()
        add_runs(bp, "•  " + b, body_size, body_color)
        bp.space_after = Pt(6)
        bp.line_spacing = 0.96
        bp._pPr.set("marL", str(Inches(0.20)))
        bp._pPr.set("indent", str(-Inches(0.20)))
    return tf


def stat(slide, x, y, w, h, value, label, fill=PANEL, value_color=NAVY):
    panel(slide, x, y, w, h, fill=fill, bar=False)
    box = slide.shapes.add_textbox(Inches(x + 0.16), Inches(y + 0.14),
                                   Inches(w - 0.32), Inches(h - 0.24))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(21)
    r.font.bold = True
    r.font.name = "Calibri"
    r.font.color.rgb = value_color
    p.space_after = Pt(2)
    lp = tf.add_paragraph()
    lr = lp.add_run()
    lr.text = label
    lr.font.size = Pt(10.5)
    lr.font.name = "Calibri"
    lr.font.color.rgb = MUTED
    lp.line_spacing = 0.95


def footnote(slide, text):
    fb = slide.shapes.add_textbox(Inches(0.55), Inches(6.94), Inches(12.23), Inches(0.32))
    tf = fb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(10.5)
    r.font.italic = True
    r.font.name = "Calibri"
    r.font.color.rgb = MUTED


def flow_box(slide, x, y, w, h, title, lines, fill=PANEL, title_color=ACCENT):
    p = slide.shapes.add_shape(RECT, Inches(x), Inches(y), Inches(w), Inches(h))
    p.fill.solid()
    p.fill.fore_color.rgb = fill
    p.line.color.rgb = RULE
    p.line.width = Pt(0.75)
    p.shadow.inherit = False
    box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.16),
                                   Inches(w - 0.28), Inches(h - 0.28))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    hp = tf.paragraphs[0]
    hr = hp.add_run()
    hr.text = title
    hr.font.size = Pt(12.5)
    hr.font.bold = True
    hr.font.name = "Calibri"
    hr.font.color.rgb = title_color
    hp.space_after = Pt(5)
    for ln in lines:
        lp = tf.add_paragraph()
        lr = lp.add_run()
        lr.text = ln
        lr.font.size = Pt(10.5)
        lr.font.name = "Calibri"
        lr.font.color.rgb = BODY
        lp.line_spacing = 0.95
        lp.space_after = Pt(2)


def arrow(slide, x, y, w=0.44, h=0.34):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False


# ---------------------------------------------------------------------------
# SLIDE 1 — problem framing (regenerated from build_problem_slide.py data)
# ---------------------------------------------------------------------------

def slide_1(prs):
    slide = new_slide(prs)
    header(slide, DECK_TITLE, PROBLEM_SUBTITLE)
    left_x, right_x = 0.55, 6.92
    top_y, bot_y = 1.62, 4.58
    w, h = 5.86, 2.78
    for idx, (heading, bullets) in enumerate(PROBLEM_SECTIONS):
        x = left_x if idx % 2 == 0 else right_x
        y = top_y if idx < 2 else bot_y
        card(slide, x, y, w, h, heading, bullets, number=idx + 1)


# ---------------------------------------------------------------------------
# SLIDE 2 — what we built (MVP)
# ---------------------------------------------------------------------------

def slide_2(prs):
    slide = new_slide(prs)
    header(slide,
           "What we built: the Transformation Roadmap Generator",
           "One browser page that turns the client's existing plans into a sequenced, "
           "defensible roadmap — no backend, no live model call.",
           eyebrow="THE MVP")

    stats = [
        ("60", "initiatives normalized"),
        ("95", "dependencies resolved"),
        ("3", "views, one click apart"),
        ("0", "runtime model calls"),
    ]
    sw, gap = 3.02, 0.055
    for i, (v, l) in enumerate(stats):
        stat(slide, 0.55 + i * (sw + gap), 1.62, sw, 0.92, v, l)

    card(slide, 0.55, 2.72, 6.13, 2.10, "WHAT IT DOES", [
        "**Ingests** 9 CSVs plus messy business cases, PMO status reports and steering minutes",
        "**Normalizes** all of it to one initiative model — cost, dates, benefit, dependencies",
        "**Sequences** with a deterministic topological sort, flagging conflicts as it goes",
    ], number=None)

    card(slide, 6.79, 2.72, 5.99, 2.10, "WHAT A JUDGE CAN CLICK", [
        "**Sequenced roadmap** — every initiative on a timeline, in dependency order",
        "**Dependency graph** — click one initiative, see everything it gates",
        "**Benefit-realization curve** — when value actually hits the P&L",
    ], number=None)

    card(slide, 0.55, 5.00, 12.23, 1.60, "WHY IT CAN'T FAIL IN THE DEMO", [
        "Static single-page app: **no backend, no runtime model call** — it renders the same way every time",
        "**AI is used at ingest only** — offline, output committed to the repo — to pull structured fields out of prose business cases",
        "**Scenario switch** re-runs the sequence in the browser; the engine is code, not a prompt",
    ], number=None)


# ---------------------------------------------------------------------------
# SLIDE 3 — the money shot
# ---------------------------------------------------------------------------

def slide_3(prs):
    slide = new_slide(prs)
    header(slide,
           "The $895k initiative ranked last — until the tool followed the chain",
           "Naive ROI buries the enabler. Dependency-aware sequencing puts it first, and shows why.",
           eyebrow="THE DEMO MOMENT")

    # before / after
    before = panel(slide, 0.55, 1.68, 5.30, 2.36, bar_color=MUTED)
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.82), Inches(4.92), Inches(2.12))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    hp = tf.paragraphs[0]
    hr = hp.add_run()
    hr.text = "RANKED BY NAIVE ROI"
    hr.font.size = Pt(13)
    hr.font.bold = True
    hr.font.name = "Calibri"
    hr.font.color.rgb = MUTED
    hp.space_after = Pt(7)
    for line in [
        "**$895k** cost, **$70k/yr** of direct benefit",
        "Scores worst in the portfolio — ranked **LAST**",
        "A PMO spreadsheet would defer it, and quietly stall everything behind it",
    ]:
        bp = tf.add_paragraph()
        add_runs(bp, "•  " + line, 12.5, BODY)
        bp.space_after = Pt(6)
        bp.line_spacing = 0.96
        bp._pPr.set("marL", str(Inches(0.20)))
        bp._pPr.set("indent", str(-Inches(0.20)))

    arrow(slide, 6.03, 2.72, 0.52, 0.40)

    card(slide, 6.73, 1.68, 6.05, 2.36, "AFTER THE TOOL FOLLOWS THE CHAIN", [
        "It walks the **28 downstream initiatives** that enabler gates",
        "Those initiatives carry **$38.6m/yr** of benefit behind it",
        "It corrects itself and sequences the enabler **first** — with the chain as the evidence",
    ], number=None, head_color=ACCENT)

    band = panel(slide, 0.55, 4.22, 12.23, 0.62, fill=RGBColor(0xEB, 0xF3, 0xF6), bar=False)
    bb = slide.shapes.add_textbox(Inches(0.75), Inches(4.34), Inches(11.83), Inches(0.42))
    btf = bb.text_frame
    btf.word_wrap = True
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    add_runs(btf.paragraphs[0],
             "**$895k of spend sits in front of $38.6m/yr of downstream benefit.** "
             "That is the insight a PMO spreadsheet cannot produce.", 14, BODY)

    card(slide, 0.55, 5.02, 12.23, 1.58, "AND IT CHALLENGES THE PLAN — 4 VIOLATION TYPES", [
        "**Dependency**, **resource**, **budget** and **compliance** violations, detected on every run",
        "**Compliance sorts to the top** — regulatory initiatives are never deferrable by a scenario",
        "Every flag is traceable to the source row, so the owner can argue with the data, not the tool",
    ], number=None)


# ---------------------------------------------------------------------------
# SLIDE 4 — business value
# ---------------------------------------------------------------------------

def slide_4(prs):
    slide = new_slide(prs)
    header(slide,
           "The board is promised $84.3m. The plan schedules a $39.6m run-rate.",
           "The gap is not a benefits problem — it is a sequencing problem, and it is now visible.",
           eyebrow="BUSINESS VALUE")

    stats = [
        ("$84.3m", "benefit promised across the portfolio"),
        ("$39.6m", "annual run-rate scheduled by end-2027 — 47% of the promise"),
        ("$18.0m", "of benefit inside the 24-month window"),
        ("$89.8m", "of planned spend inside that same window"),
    ]
    sw, gap = 3.02, 0.055
    for i, (v, l) in enumerate(stats):
        stat(slide, 0.55 + i * (sw + gap), 1.62, sw, 1.10, v, l)

    card(slide, 0.55, 2.92, 6.13, 1.94, "WHAT THE NUMBERS SAY", [
        "**16 of 60** initiatives deliver nothing inside the window",
        "Only **25 of 60** ever reach full run-rate",
        "Budget: **$95.3m** approved against a **$101.7m** forecast",
    ], number=None)

    card(slide, 6.79, 2.92, 5.99, 1.94, "WHAT THAT IS WORTH", [
        "**Weeks of PMO analysis compressed to minutes**, repeatable on every status cycle",
        "A board that can finally see **when value hits the P&L** — not just what was promised",
        "Re-phasing decisions get made **before** money is committed, not after",
    ], number=None)

    card(slide, 0.55, 5.04, 12.23, 1.56, "THE CONVERSATION IT UNLOCKS", [
        "\"Which initiatives are we funding this year that return nothing inside the window?\" — answered in one view",
        "\"What has to move for benefit to land before the end of 2027?\" — answered by re-running the sequence",
        "Same data the client already has. The tool just refuses to let the timing hide.",
    ], number=None)


# ---------------------------------------------------------------------------
# SLIDE 5 — path to market (highest-weight)
# ---------------------------------------------------------------------------

def slide_5(prs):
    slide = new_slide(prs)
    header(slide,
           "Path to market: land it in a CTO's first 90 days",
           "One client portfolio, 4-6 weeks, their real PMO exports — out the other end, a board-ready roadmap.",
           eyebrow="PATH TO MARKET")

    card(slide, 0.55, 1.62, 4.02, 2.42, "WHO BUYS IT", [
        "**Primary: Chief Transformation Officer**, first 90 days, owes the board a credible plan",
        "**Secondary: PMO lead** — maintains the roadmap between board cycles",
        "**Secondary: CFO** — defends the benefit baseline and the phasing",
    ], number=None)

    card(slide, 4.66, 1.62, 4.02, 2.42, "THE PILOT", [
        "**One client transformation portfolio, 4-6 weeks**",
        "Ingest **their real PMO exports** — trackers, business cases, status reports",
        "Deliver a **board-ready sequenced roadmap** plus the conflicts we found",
    ], number=None)

    card(slide, 8.77, 1.62, 4.01, 2.42, "WHAT WE NEED FROM THEM", [
        "A **read-only export** of their programme tracker",
        "**One PMO analyst, part-time**, to confirm what the data means",
        "A **named executive sponsor** who will act on the sequencing",
    ], number=None)

    card(slide, 0.55, 4.22, 8.13, 2.38, "RISKS AND HOW WE HANDLE THEM", [
        "**Data quality** — the tool **surfaces contradictions rather than hiding them**; the mess becomes the first deliverable",
        "**Change resistance from initiative owners** — sequencing is **transparent and auditable**, every ordering traces to a dependency, not a black box",
        "**Security** — runs **client-side**, no data leaves their environment, nothing is uploaded and no model is called at runtime",
    ], number=None)

    card(slide, 8.77, 4.22, 4.01, 2.38, "HOW WE MEASURE SUCCESS", [
        "**Time to produce a board roadmap** — weeks to minutes",
        "**Sequencing conflicts caught pre-commitment**",
        "**% of benefit re-phased into the window**",
    ], number=None, bar_color=NAVY, head_color=NAVY)


# ---------------------------------------------------------------------------
# SLIDE 6 — how the dashboard is used
# ---------------------------------------------------------------------------

def slide_6(prs):
    slide = new_slide(prs)
    header(slide,
           "How the dashboard is used",
           "Three things a Chief Transformation Officer does with it in a single sitting.",
           eyebrow="USING THE TOOL")

    y, h, w = 1.68, 3.28, 4.02

    card(slide, 0.55, y, w, h, "1.  EVALUATING TRADEOFFS", [
        "Switch between **board scenarios** with one control",
        "The **sequence**, the **conflicts** and the **benefit curve** all move together, side by side",
        "Compare **\"fastest value\"** against **\"lowest risk\"** on screen",
        "The choice gets **compared, not argued about**",
    ], number=None)

    card(slide, 4.66, y, w, h, "2.  PRIORITIZING AGAINST CONSTRAINTS", [
        "Sequences against the **real caps in the data** — budget and capex caps, peak FTE caps, mandatory and deferred initiatives",
        "Flags **4 violation types**: dependency, resource, budget, compliance",
        "**Compliance sorts to the top** of the exception list",
        "**Regulatory initiatives can never be deferred** by a scenario",
    ], number=None)

    card(slide, 8.77, y, 4.01, h, "3.  VALUE DELIVERED TO THE CLIENT", [
        "A **defensible sequencing decision in minutes**, not weeks of PMO analysis",
        "**Conflicts caught before money is committed**, not after",
        "A clear answer for the board on **when benefit reaches the P&L**",
    ], number=None, bar_color=NAVY, head_color=NAVY)

    band = panel(slide, 0.55, 5.24, 12.23, 1.36, fill=RGBColor(0xEB, 0xF3, 0xF6), bar=False)
    bb = slide.shapes.add_textbox(Inches(0.80), Inches(5.44), Inches(11.73), Inches(1.00))
    btf = bb.text_frame
    btf.word_wrap = True
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    btf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate([
        "**One screen, one owner, one monthly cycle.** The CTO drives it in the board meeting; "
        "the PMO lead re-runs it on every status report.",
        "Nothing is hidden behind a model call — every ordering and every exception traces back to a row in the client's own data.",
    ]):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        add_runs(p, line, 13, BODY)
        p.space_after = Pt(5)
        p.line_spacing = 0.98


# ---------------------------------------------------------------------------
# SLIDE 7 — end-to-end workflow / architecture
# ---------------------------------------------------------------------------

def slide_7(prs):
    slide = new_slide(prs)
    header(slide,
           "The workflow, end to end: from ingestion to output",
           "AI at ingest, deterministic code at decision time — re-run it and a judge gets the same answer, every time.",
           eyebrow="WORKFLOW: INGESTION → OUTPUT")

    y, h = 2.00, 2.34
    xs = [0.55, 3.72, 6.89, 10.06]
    w = 2.72

    # step ribbon above the flow, so the left-to-right reading order is explicit
    labels = ["STEP 1  ·  INGEST", "STEP 2  ·  EXTRACT", "STEP 3  ·  NORMALIZE", "STEP 4  ·  OUTPUT"]
    for i, lab in enumerate(labels):
        rb = slide.shapes.add_shape(RECT, Inches(xs[i]), Inches(1.62), Inches(w if i < 3 else 2.72), Inches(0.30))
        rb.fill.solid()
        rb.fill.fore_color.rgb = ACCENT
        rb.line.fill.background()
        rb.shadow.inherit = False
        tb = slide.shapes.add_textbox(Inches(xs[i]), Inches(1.665), Inches(w), Inches(0.24))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = lab
        r.font.size = Pt(10.5)
        r.font.bold = True
        r.font.name = "Calibri"
        r.font.color.rgb = WHITE

    flow_box(slide, xs[0], y, w, h, "WHAT EXISTS TODAY", [
        "9 CSV extracts",
        "Prose business cases",
        "PMO status reports",
        "Steering minutes",
        "A messy master tracker",
    ])
    arrow(slide, xs[0] + w + 0.10, y + h / 2 - 0.17)

    flow_box(slide, xs[1], y, w, h, "OFFLINE INGEST", [
        "AI extraction of structured",
        "fields from prose",
        "+ normalization rules",
        "",
        "Run once, output committed",
    ])
    arrow(slide, xs[1] + w + 0.10, y + h / 2 - 0.17)

    flow_box(slide, xs[2], y, w, h, "ONE DATA FILE", [
        "A single committed",
        "initiative model:",
        "cost, dates, benefit,",
        "dependencies, constraints",
        "",
        "The app's only input",
    ], fill=RGBColor(0xEB, 0xF3, 0xF6))
    arrow(slide, xs[2] + w + 0.10, y + h / 2 - 0.17)

    flow_box(slide, xs[3], y, w, h, "IN THE BROWSER", [
        "Deterministic engine:",
        "topological sort +",
        "conflict detection",
        "",
        "→ Sequenced roadmap",
        "→ Dependency graph",
        "→ Benefit curve",
    ])

    card(slide, 0.55, 4.62, 12.23, 1.72, "WHY THIS SHAPE", [
        "**Sequencing is deterministic code, not a model call** — the same input always produces the same roadmap",
        "**AI does the one job it is good at**: turning prose business cases into structured fields, offline, reviewable in the diff",
        "**No backend, no keys, no network at runtime** — it runs from a file, which is also why it is safe inside a client's environment",
    ], number=None)


# ---------------------------------------------------------------------------
# SLIDE 8 — Aberdeen Labs reusability
# ---------------------------------------------------------------------------

def slide_8(prs):
    slide = new_slide(prs)
    header(slide,
           "A Labs accelerator, not a one-off demo",
           "Three separable layers — the next team starts from a working dataset, not a blank repo.",
           eyebrow="ABERDEEN LABS REUSABILITY")

    y, h, w = 1.68, 2.30, 4.02
    flow_box(slide, 0.55, y, w, h, "INGEST → MODEL", [
        "Extraction + normalization",
        "into one initiative schema.",
        "",
        "Swap the source documents,",
        "keep everything downstream.",
    ])
    flow_box(slide, 4.66, y, w, h, "SEQUENCING ENGINE", [
        "Topological sort +",
        "4-type conflict detection.",
        "",
        "Domain-agnostic: it only",
        "needs initiatives and edges.",
    ], fill=RGBColor(0xEB, 0xF3, 0xF6))
    flow_box(slide, 8.77, y, 4.01, h, "VIEW LAYER", [
        "Roadmap, dependency graph,",
        "benefit curve, scenarios.",
        "",
        "Re-skinnable; reads the",
        "model, not the source data.",
    ])

    card(slide, 0.55, 4.16, 6.13, 2.44, "THE SAME ENGINE, OTHER ENGAGEMENTS", [
        "**Post-merger integration** — synergy programmes with hard sequencing constraints",
        "**Cost programmes** — savings initiatives that gate each other",
        "**Capital plans** — capex portfolios where enablers come first",
        "Any **dependency-heavy portfolio** where timing decides value",
    ], number=None)

    card(slide, 6.79, 4.16, 5.99, 2.44, "ALREADY IN THE REPO", [
        "**Schema** and **FIELD-REFERENCE.md** for the initiative model",
        "**Synthetic test portfolio** — 60 initiatives, 95 dependencies, deliberate contradictions",
        "**Validator** that checks a new dataset before it reaches the engine",
        "Next team starts from a **working dataset**, day one",
    ], number=None, bar_color=NAVY, head_color=NAVY)


# ---------------------------------------------------------------------------

BUILDERS = [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8]


def build(path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in BUILDERS:
        fn(prs)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"wrote {path} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    here = os.path.dirname(os.path.abspath(__file__))
    build(os.path.join(here, "transformation-roadmap-generator.pptx"))
